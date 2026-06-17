"""PPTX Text Extraction module for IEEE Content Conversion pipeline.

Extracts slide text (shapes, tables, notes) from .pptx files stored in S3
using python-pptx. Returns cleaned text suitable for passing to Bedrock.
"""

from __future__ import annotations

import json
import logging
import re
from datetime import datetime, timezone
from io import BytesIO
from typing import TypedDict

import boto3
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError

logger = logging.getLogger(__name__)

MAX_TEXT_LENGTH = 180_000


class ExtractionResult(TypedDict):
    text: str
    slide_count: int
    extraction_method: str  # "extract_text" | "failed"


class PPTXExtractor:
    """Extracts text from .pptx presentations stored in S3.

    Usage (called by the orchestrator):
        extractor = PPTXExtractor(s3_client=boto3.client("s3"))
        result = extractor.extract(bucket, key, ou, product_part_number)
    """

    def __init__(self, s3_client=None):
        self._s3 = s3_client or boto3.client("s3")

    def extract(
        self,
        bucket: str,
        key: str,
        ou: str,
        product_part_number: str,
    ) -> ExtractionResult:
        """Download a .pptx from S3, extract text, and write metadata.

        Args:
            bucket: S3 bucket name.
            key: S3 object key, e.g. ``{ou}/pending/{filename}.pptx``.
            ou: Organizational unit prefix used for output paths.
            product_part_number: Identifier used in the metadata filename.

        Returns:
            ExtractionResult with text, slide_count, and extraction_method.
        """
        logger.info("Extracting text from s3://%s/%s", bucket, key)
        pptx_bytes = self._download(bucket, key)
        result = self.extract_from_bytes(pptx_bytes)

        self._write_metadata(
            bucket=bucket,
            ou=ou,
            product_part_number=product_part_number,
            slide_count=result["slide_count"],
            extraction_method=result["extraction_method"],
        )

        return result

    def extract_from_bytes(self, pptx_bytes: bytes) -> ExtractionResult:
        """Extract text from raw .pptx bytes (no S3 interaction).

        Useful for unit testing and non-S3 callers.
        """
        try:
            prs = Presentation(BytesIO(pptx_bytes))
        except PackageNotFoundError:
            logger.exception("Invalid or corrupt .pptx file")
            return ExtractionResult(text="", slide_count=0, extraction_method="failed")
        except Exception:
            logger.exception("Failed to open .pptx")
            return ExtractionResult(text="", slide_count=0, extraction_method="failed")

        slides = list(prs.slides)
        slide_count = len(slides)

        slide_texts: list[str] = []
        has_text = False

        for slide in slides:
            slide_text = _extract_slide_text(slide)
            if slide_text.strip():
                has_text = True
            slide_texts.append(slide_text)

        if not has_text:
            logger.warning(
                "No extractable text found across %d slides", slide_count
            )
            return ExtractionResult(
                text="", slide_count=slide_count, extraction_method="extract_text"
            )

        full_text = "\n\n".join(slide_texts)
        full_text = _clean_text(full_text)

        if len(full_text) > MAX_TEXT_LENGTH:
            logger.info(
                "Truncating extracted text from %d to %d characters",
                len(full_text),
                MAX_TEXT_LENGTH,
            )
            full_text = full_text[:MAX_TEXT_LENGTH]

        return ExtractionResult(
            text=full_text, slide_count=slide_count, extraction_method="extract_text"
        )

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _download(self, bucket: str, key: str) -> bytes:
        resp = self._s3.get_object(Bucket=bucket, Key=key)
        return resp["Body"].read()

    def _write_metadata(
        self,
        bucket: str,
        ou: str,
        product_part_number: str,
        slide_count: int,
        extraction_method: str,
    ) -> None:
        metadata_key = f"{ou}/metadata/{product_part_number}.pptx.json"
        body = json.dumps(
            {
                "slideCount": slide_count,
                "extractionMethod": extraction_method,
                "extractedAt": datetime.now(timezone.utc)
                .isoformat()
                .replace("+00:00", "Z"),
            }
        )
        self._s3.put_object(
            Bucket=bucket,
            Key=metadata_key,
            Body=body.encode(),
            ContentType="application/json",
        )
        logger.info("Wrote metadata to s3://%s/%s", bucket, metadata_key)


# ----------------------------------------------------------------------
# Slide text extraction
# ----------------------------------------------------------------------

def _extract_slide_text(slide) -> str:
    """Collect text from all shapes, tables, and notes on a slide."""
    fragments: list[str] = []

    for shape in slide.shapes:
        fragments.extend(_extract_shape_text(shape))

    notes = getattr(slide, "notes_slide", None)
    if notes is not None and notes.notes_text_frame is not None:
        notes_text = notes.notes_text_frame.text.strip()
        if notes_text:
            fragments.append(notes_text)

    return "\n".join(f for f in fragments if f)


def _extract_shape_text(shape) -> list[str]:
    """Extract text from a single shape, including grouped shapes and tables."""
    fragments: list[str] = []

    # Grouped shapes — recurse.
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for inner in shape.shapes:
                fragments.extend(_extract_shape_text(inner))
        except AttributeError:
            pass
        return fragments

    # Tables — collect per-cell text.
    if getattr(shape, "has_table", False):
        try:
            table = shape.table
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells]
                row_text = " | ".join(c for c in row_cells if c)
                if row_text:
                    fragments.append(row_text)
        except AttributeError:
            pass
        return fragments

    # Normal text shapes.
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text.strip()
        if text:
            fragments.append(text)

    return fragments


# ----------------------------------------------------------------------
# Text cleaning utilities
# ----------------------------------------------------------------------

_EXCESSIVE_NEWLINES_RE = re.compile(r"\n{3,}")


def _clean_text(text: str) -> str:
    """Normalise whitespace in extracted slide text."""
    text = _EXCESSIVE_NEWLINES_RE.sub("\n\n", text)
    return text.strip()
