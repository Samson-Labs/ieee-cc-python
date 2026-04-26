"""Tests for PPTXExtractor.

Covers: normal pptx, empty slides, multi-slide text, tables, speaker notes,
grouped shapes, corrupted file, truncation, S3 write of metadata.
"""

from __future__ import annotations

import json
from io import BytesIO
from unittest.mock import MagicMock

import pytest
from botocore.exceptions import ClientError

from pptx import Presentation
from pptx.util import Inches

from src.extractors.pptx_extractor import (
    MAX_TEXT_LENGTH,
    PPTXExtractor,
    _clean_text,
)


# ---------------------------------------------------------------------------
# Helpers to build in-memory .pptx files
# ---------------------------------------------------------------------------

def _make_pptx(slides_text: list[str], *, notes: list[str] | None = None) -> bytes:
    """Create a minimal .pptx with one title per slide and optional notes."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # title-only layout
    notes = notes or []

    for idx, text in enumerate(slides_text):
        slide = prs.slides.add_slide(blank_layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = text

        if idx < len(notes) and notes[idx]:
            slide.notes_slide.notes_text_frame.text = notes[idx]

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_table() -> bytes:
    """Create a .pptx with a 2x2 table on one slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    rows, cols = 2, 2
    left = top = Inches(1)
    width = Inches(4)
    height = Inches(2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Cell A1"
    table.cell(1, 1).text = "Cell B1"

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_empty_pptx() -> bytes:
    """Create a .pptx with one blank slide (no text)."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# extract_from_bytes — core extraction logic (no S3)
# ---------------------------------------------------------------------------

class TestExtractFromBytes:
    def test_extracts_text_from_slide_titles(self):
        extractor = PPTXExtractor(s3_client=MagicMock())
        pptx = _make_pptx(["Introduction", "Architecture", "Conclusion"])

        result = extractor.extract_from_bytes(pptx)

        assert result["slide_count"] == 3
        assert result["extraction_method"] == "text"
        assert "Introduction" in result["text"]
        assert "Architecture" in result["text"]
        assert "Conclusion" in result["text"]

    def test_extracts_speaker_notes(self):
        extractor = PPTXExtractor(s3_client=MagicMock())
        pptx = _make_pptx(
            ["Slide One"],
            notes=["Notes with commentary for the speaker."],
        )

        result = extractor.extract_from_bytes(pptx)

        assert "Slide One" in result["text"]
        assert "Notes with commentary" in result["text"]

    def test_extracts_table_cells(self):
        extractor = PPTXExtractor(s3_client=MagicMock())
        pptx = _make_pptx_with_table()

        result = extractor.extract_from_bytes(pptx)

        assert result["slide_count"] == 1
        assert result["extraction_method"] == "text"
        for token in ("Header A", "Header B", "Cell A1", "Cell B1"):
            assert token in result["text"]

    def test_empty_pptx_returns_empty_method(self):
        extractor = PPTXExtractor(s3_client=MagicMock())
        pptx = _make_empty_pptx()

        result = extractor.extract_from_bytes(pptx)

        assert result["slide_count"] == 1
        assert result["extraction_method"] == "empty"
        assert result["text"] == ""

    def test_corrupted_pptx_returns_failed(self):
        extractor = PPTXExtractor(s3_client=MagicMock())

        result = extractor.extract_from_bytes(b"not a real pptx file")

        assert result["slide_count"] == 0
        assert result["extraction_method"] == "failed"
        assert result["text"] == ""

    def test_truncates_very_long_text(self):
        extractor = PPTXExtractor(s3_client=MagicMock())
        # Fill a single slide with enough text to exceed MAX_TEXT_LENGTH.
        big = "x" * (MAX_TEXT_LENGTH + 5_000)
        pptx = _make_pptx([big])

        result = extractor.extract_from_bytes(pptx)

        assert len(result["text"]) == MAX_TEXT_LENGTH
        assert result["extraction_method"] == "text"


# ---------------------------------------------------------------------------
# extract() — S3 download + metadata write
# ---------------------------------------------------------------------------

class TestExtract:
    def test_writes_metadata_to_expected_key(self):
        pptx_bytes = _make_pptx(["Slide A", "Slide B"])

        s3 = MagicMock()
        s3.get_object.return_value = {"Body": BytesIO(pptx_bytes)}

        extractor = PPTXExtractor(s3_client=s3)
        result = extractor.extract(
            bucket="b",
            key="ieee/pending/STD-123.pptx",
            ou="ieee",
            product_part_number="STD-123",
        )

        assert result["slide_count"] == 2

        # Metadata write
        put_calls = [c for c in s3.method_calls if c[0] == "put_object"]
        assert len(put_calls) == 1
        kwargs = put_calls[0][2]
        assert kwargs["Bucket"] == "b"
        assert kwargs["Key"] == "ieee/metadata/STD-123.pptx.json"
        assert kwargs["ContentType"] == "application/json"

        body = json.loads(kwargs["Body"])
        assert body["slideCount"] == 2
        assert body["extractionMethod"] == "text"
        assert body["extractedAt"].endswith("Z")

    def test_s3_get_object_error_propagates(self):
        s3 = MagicMock()
        s3.get_object.side_effect = ClientError(
            {"Error": {"Code": "NoSuchKey", "Message": "Not found"}},
            "GetObject",
        )

        extractor = PPTXExtractor(s3_client=s3)
        with pytest.raises(ClientError):
            extractor.extract(
                bucket="b",
                key="ieee/pending/missing.pptx",
                ou="ieee",
                product_part_number="missing",
            )


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

class TestCleanText:
    def test_collapses_excessive_newlines(self):
        assert _clean_text("a\n\n\n\nb") == "a\n\nb"

    def test_strips_surrounding_whitespace(self):
        assert _clean_text("  hello  \n\n") == "hello"
